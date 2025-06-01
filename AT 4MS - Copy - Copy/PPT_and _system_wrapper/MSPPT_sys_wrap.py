"""
PowerPoint Wrapper for Gestures control
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image
import win32com.client

class PowerPointWrapper:
    """A wrapper for PowerPoint file operations and live PowerPoint control."""

    def __init__(self, pptx_path=None):
        """
        Initialize the wrapper. If pptx_path is given, open or create the file.
        """
        self.pptx_path = pptx_path
        if pptx_path and os.path.exists(pptx_path):
            self.prs = Presentation(pptx_path)
        else:
            self.prs = Presentation()
        self.app = None  # For live PowerPoint control

    def add_title_slide(self, title, subtitle=""):
        """Add a title slide."""
        slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = subtitle

    def add_text_slide(self, title, content):
        """Add a slide with a title and text content."""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = content

    def add_image_slide(self, title, image_path):
        """Add a slide with a title and an image."""
        slide_layout = self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        left = Inches(1)
        top = Inches(2)
        slide.shapes.add_picture(image_path, left, top, width=Inches(5))

    def save(self, path=None):
        """Save the presentation."""
        save_path = path or self.pptx_path or "output.pptx"
        self.prs.save(save_path)
        print(f"Presentation saved to {save_path}")

    def open_in_powerpoint(self, path=None):
        """Open the presentation in Microsoft PowerPoint (requires Windows)."""
        ppt_path = os.path.abspath(path or self.pptx_path or "output.pptx")
        self.app = win32com.client.Dispatch("PowerPoint.Application")
        self.app.Visible = True
        self.app.Presentations.Open(ppt_path)
        print(f"Opened {ppt_path} in PowerPoint.")

    def close_powerpoint(self):
        """Close PowerPoint application."""
        if self.app:
            self.app.Quit()
            self.app = None

    def add_prediction_slide(self, image_path, prediction, class_names=None):
        """Add a slide showing the input image and prediction result."""
        slide_layout = self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Model Prediction"
        left = Inches(0.5)
        top = Inches(1.5)
        slide.shapes.add_picture(image_path, left, top, width=Inches(3))
        txBox = slide.shapes.add_textbox(Inches(4), Inches(1.5), Inches(4), Inches(2))
        tf = txBox.text_frame
        if class_names:
            pred_label = class_names[int(prediction)]
            tf.text = f"Prediction: {pred_label}"
        else:
            tf.text = f"Prediction: {prediction}"
        tf.paragraphs[0].font.size = Pt(24)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(0, 0, 255)
        tf.paragraphs[0].alignment = PP_ALIGN.LEFT

def send_result_to_powerpoint(image_path, prediction, class_names=None, pptx_path="output.pptx"):
    """
    Example function to send model results to PowerPoint.
    """
    ppt = PowerPointWrapper(pptx_path)
    ppt.add_prediction_slide(image_path, prediction, class_names)
    ppt.save(pptx_path)
    ppt.open_in_powerpoint(pptx_path)
